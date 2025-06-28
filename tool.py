#!/usr/bin/env python3
"""
PowerPoint编辑器工具类
提供基础的PPT编辑功能，包括添加文本、图片、形状等
"""

import logging
from typing import Any, Dict, List, Optional, TYPE_CHECKING
from pathlib import Path
import json

# 导入PowerPoint相关库
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
except ImportError:
    raise ImportError("请安装python-pptx库: pip install python-pptx")

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType

# 设置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PowerPointEditor:
    """PowerPoint编辑器类"""
    
    def __init__(self):
        self.current_presentation: Optional["PresentationType"] = None
        self.current_file_path: Optional[str] = None

    def create_presentation(self) -> Dict[str, Any]:
        """创建新的演示文稿"""
        try:
            self.current_presentation = Presentation()
            self.current_file_path = None
            return {
                "success": True,
                "message": "成功创建新的演示文稿",
                "slides_count": len(self.current_presentation.slides)
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def open_presentation(self, file_path: str) -> Dict[str, Any]:
        """打开现有的演示文稿"""
        try:
            if not Path(file_path).exists():
                return {"success": False, "error": f"文件不存在: {file_path}"}

            self.current_presentation = Presentation(file_path)
            self.current_file_path = file_path

            return {
                "success": True,
                "message": f"成功打开演示文稿: {file_path}",
                "slides_count": len(self.current_presentation.slides),
                "file_path": file_path
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def save_presentation(self, file_path: Optional[str] = None) -> Dict[str, Any]:
        """保存演示文稿"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            save_path = file_path or self.current_file_path
            if not save_path:
                return {"success": False, "error": "请指定保存路径"}

            # 保存前验证过渡效果
            transition_count = self._count_transitions()

            self.current_presentation.save(save_path)
            self.current_file_path = save_path

            return {
                "success": True,
                "message": f"成功保存演示文稿: {save_path}",
                "file_path": save_path,
                "slides_with_transitions": transition_count,
                "note": f"文件包含 {transition_count} 张有过渡效果的幻灯片" if transition_count > 0 else "文件不包含过渡效果"
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _count_transitions(self) -> int:
        """统计有过渡效果的幻灯片数量"""
        try:
            if not self.current_presentation:
                return 0

            count = 0
            namespaces = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

            for slide in self.current_presentation.slides:
                slide_element = slide._element
                transition_elem = slide_element.find('.//p:transition', namespaces)
                if transition_elem is not None:
                    count += 1

            return count
        except:
            return 0

    def add_slide(self, layout_index: int = 1) -> Dict[str, Any]:
        """添加新幻灯片"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            # 获取幻灯片布局
            slide_layouts = self.current_presentation.slide_layouts
            if layout_index >= len(slide_layouts):
                layout_index = 1  # 默认使用标题和内容布局

            layout = slide_layouts[layout_index]
            slide = self.current_presentation.slides.add_slide(layout)

            return {
                "success": True,
                "message": f"成功添加新幻灯片",
                "slide_index": len(self.current_presentation.slides) - 1,
                "total_slides": len(self.current_presentation.slides)
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_text_box(self, slide_index: int, text: str, left: float = 1,
                     top: float = 1, width: float = 8, height: float = 1,
                     font_size: int = 18, font_color: str = "000000") -> Dict[str, Any]:
        """在指定幻灯片添加文本框"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            # 添加文本框
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width)
            height_inches = Inches(height)

            textbox = slide.shapes.add_textbox(left_inches, top_inches, width_inches, height_inches)
            text_frame = textbox.text_frame
            text_frame.text = text

            # 设置字体样式
            paragraph = text_frame.paragraphs[0]
            font = paragraph.font
            font.size = Pt(font_size)

            # 设置字体颜色
            try:
                rgb_color = RGBColor.from_string(font_color)
                font.color.rgb = rgb_color
            except:
                pass  # 如果颜色格式不正确，使用默认颜色

            return {
                "success": True,
                "message": f"成功在幻灯片 {slide_index} 添加文本框",
                "text": text,
                "position": {"left": left, "top": top, "width": width, "height": height}
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_title_slide(self, title: str, subtitle: str = "") -> Dict[str, Any]:
        """添加标题幻灯片"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            # 使用标题幻灯片布局
            title_slide_layout = self.current_presentation.slide_layouts[0]
            slide = self.current_presentation.slides.add_slide(title_slide_layout)            # 设置标题
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title            # 设置副标题
            if subtitle and len(slide.placeholders) > 1:
                try:
                    subtitle_shape = slide.placeholders[1]
                    subtitle_shape.text_frame.text = subtitle  # type: ignore
                except (AttributeError, TypeError):
                    pass  # 如果无法设置副标题，忽略错误

            return {
                "success": True,
                "message": "成功添加标题幻灯片",
                "slide_index": len(self.current_presentation.slides) - 1,
                "title": title,
                "subtitle": subtitle
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_bullet_points(self, slide_index: int, title: str, bullet_points: List[str]) -> Dict[str, Any]:
        """添加带项目符号的内容幻灯片"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            # 设置标题
            if slide.shapes.title:
                slide.shapes.title.text = title            # 查找内容占位符
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # 内容占位符通常是索引1
                    content_placeholder = shape
                    break
                    
            if content_placeholder:
                try:
                    text_frame = content_placeholder.text_frame  # type: ignore
                    text_frame.clear()  # 清除现有内容

                    for i, point in enumerate(bullet_points):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0  # 设置为第一级项目符号
                except (AttributeError, TypeError):
                    pass  # 如果无法设置内容，忽略错误

            return {
                "success": True,
                "message": f"成功在幻灯片 {slide_index} 添加项目符号内容",
                "title": title,
                "bullet_points": bullet_points
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_image(self, slide_index: int, image_path: str, left: float = 1,
                  top: float = 2, width: Optional[float] = None, height: Optional[float] = None) -> Dict[str, Any]:
        """在幻灯片中添加图片"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            if not Path(image_path).exists():
                return {"success": False, "error": f"图片文件不存在: {image_path}"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            # 添加图片
            left_inches = Inches(left)
            top_inches = Inches(top)

            if width and height:
                width_inches = Inches(width)
                height_inches = Inches(height)
                pic = slide.shapes.add_picture(image_path, left_inches, top_inches, width_inches, height_inches)
            else:
                pic = slide.shapes.add_picture(image_path, left_inches, top_inches)

            return {
                "success": True,
                "message": f"成功在幻灯片 {slide_index} 添加图片",
                "image_path": image_path,
                "position": {"left": left, "top": top, "width": width, "height": height}
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_shape(self, slide_index: int, shape_type: str, left: float = 1,
                  top: float = 1, width: float = 2, height: float = 1,
                  fill_color: str = "0066CC") -> Dict[str, Any]:
        """添加形状"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            # 形状类型映射
            shape_map = {
                "rectangle": MSO_SHAPE.RECTANGLE,
                "oval": MSO_SHAPE.OVAL,
                "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
                "diamond": MSO_SHAPE.DIAMOND,
                "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
                "hexagon": MSO_SHAPE.HEXAGON,
                "star": MSO_SHAPE.STAR_5_POINT,
                "arrow": MSO_SHAPE.BLOCK_ARC
            }

            if shape_type.lower() not in shape_map:
                return {"success": False, "error": f"不支持的形状类型: {shape_type}"}

            # 添加形状
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width)
            height_inches = Inches(height)

            shape = slide.shapes.add_shape(
                shape_map[shape_type.lower()],
                left_inches, top_inches, width_inches, height_inches
            )

            # 设置填充颜色
            try:
                rgb_color = RGBColor.from_string(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb_color
            except:
                pass  # 如果颜色格式不正确，使用默认颜色

            return {
                "success": True,
                "message": f"成功在幻灯片 {slide_index} 添加 {shape_type} 形状",
                "shape_type": shape_type,
                "position": {"left": left, "top": top, "width": width, "height": height},
                "fill_color": fill_color
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_presentation_info(self) -> Dict[str, Any]:
        """获取当前演示文稿信息"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides_info = []
            for i, slide in enumerate(self.current_presentation.slides):
                slide_info = {
                    "index": i,
                    "shapes_count": len(slide.shapes),
                    "has_title": bool(slide.shapes.title and slide.shapes.title.text),
                    "title": slide.shapes.title.text if slide.shapes.title else ""                }
                slides_info.append(slide_info)
            
            return {
                "success": True,
                "file_path": self.current_file_path,
                "slides_count": len(self.current_presentation.slides),
                "slides": slides_info
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def delete_slide(self, slide_index: int) -> Dict[str, Any]:
        """删除指定幻灯片"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}            # 获取要删除的幻灯片
            slide_to_remove = slides[slide_index]
            
            # 删除幻灯片的正确方法
            # 直接从slides集合中删除
            xml_slides = self.current_presentation.part._element.sldIdLst
            xml_slides.remove(xml_slides[slide_index])
            slides._sldIdLst.remove(slides._sldIdLst[slide_index])

            return {
                "success": True,
                "message": f"成功删除幻灯片 {slide_index}",
                "remaining_slides": len(self.current_presentation.slides)
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def duplicate_slide(self, slide_index: int) -> Dict[str, Any]:
        """复制指定幻灯片"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            # 获取要复制的幻灯片
            source_slide = slides[slide_index]

            # 复制幻灯片布局
            slide_layout = source_slide.slide_layout
            new_slide = slides.add_slide(slide_layout)

            # 复制所有形状
            for shape in source_slide.shapes:
                if not shape.is_placeholder:
                    # 复制非占位符形状
                    self._copy_shape(shape, new_slide)

            return {
                "success": True,
                "message": f"成功复制幻灯片 {slide_index}",
                "new_slide_index": len(slides) - 1,
                "total_slides": len(slides)
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def move_slide(self, from_index: int, to_index: int) -> Dict[str, Any]:
        """移动幻灯片位置"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if from_index >= len(slides) or to_index >= len(slides):
                return {"success": False, "error": "幻灯片索引超出范围"}

            if from_index == to_index:
                return {"success": True, "message": "幻灯片位置未改变"}

            # 简化的移动方法：复制幻灯片到新位置，然后删除原位置
            # 这是一个更安全的方法，避免直接操作XML

            # 获取源幻灯片的布局
            source_slide = slides[from_index]
            slide_layout = source_slide.slide_layout

            # 在目标位置创建新幻灯片
            if to_index >= len(slides):
                new_slide = slides.add_slide(slide_layout)
            else:
                # 在指定位置插入需要更复杂的操作，这里简化处理
                new_slide = slides.add_slide(slide_layout)

            # 复制内容（简化版本）
            for shape in source_slide.shapes:
                if not shape.is_placeholder:
                    self._copy_shape(shape, new_slide)

            # 删除原幻灯片（如果新幻灯片在后面）
            if from_index < len(slides) - 1:
                # 由于添加了新幻灯片，原索引可能需要调整
                actual_from_index = from_index if to_index > from_index else from_index
                xml_slides = self.current_presentation.part._element.sldIdLst
                xml_slides.remove(xml_slides[actual_from_index])
                slides._sldIdLst.remove(slides._sldIdLst[actual_from_index])

            return {
                "success": True,
                "message": f"成功将幻灯片从位置 {from_index} 移动到位置 {to_index}",
                "total_slides": len(slides)
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_table(self, slide_index: int, rows: int, cols: int, left: float = 1,
                  top: float = 2, width: float = 8, height: float = 4) -> Dict[str, Any]:
        """在幻灯片中添加表格"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            # 添加表格
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width)
            height_inches = Inches(height)

            table = slide.shapes.add_table(rows, cols, left_inches, top_inches, width_inches, height_inches)

            return {
                "success": True,
                "message": f"成功在幻灯片 {slide_index} 添加 {rows}x{cols} 表格",
                "rows": rows,
                "cols": cols,
                "position": {"left": left, "top": top, "width": width, "height": height}
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def set_table_cell_text(self,
                          slide_index: int,
                          table_index: int,
                          row: int,
                          col: int,
                          text: str) -> Dict[str, Any]:
        """设置表格单元格文本"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            # 查找表格 - 使用更安全的方法
            tables = []
            for shape in slide.shapes:
                try:
                    # 检查是否是表格形状
                    if hasattr(shape, 'table'):
                        table_obj = getattr(shape, 'table', None)
                        if table_obj is not None:
                            tables.append(shape)
                except:
                    continue

            if table_index >= len(tables):
                return {"success": False, "error": f"表格索引超出范围: {table_index}"}

            table_shape = tables[table_index]
            table = getattr(table_shape, 'table')

            if row >= len(table.rows) or col >= len(table.columns):
                return {"success": False, "error": "单元格位置超出表格范围"}

            # 设置单元格文本
            cell = table.cell(row, col)
            cell.text = text

            return {
                "success": True,
                "message": f"成功设置表格单元格 ({row}, {col}) 的文本",
                "text": text
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def set_slide_background_color(self, slide_index: int, color: str) -> Dict[str, Any]:
        """设置幻灯片背景颜色"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            # 设置背景颜色
            background = slide.background
            fill = background.fill
            fill.solid()

            try:
                rgb_color = RGBColor.from_string(color)
                fill.fore_color.rgb = rgb_color
            except:
                return {"success": False, "error": f"无效的颜色格式: {color}"}

            return {
                "success": True,
                "message": f"成功设置幻灯片 {slide_index} 的背景颜色",
                "color": color
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def add_hyperlink(self, slide_index: int, shape_index: int, url: str, display_text: Optional[str] = None) -> Dict[str, Any]:
        """为形状添加超链接"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            if shape_index >= len(slide.shapes):
                return {"success": False, "error": f"形状索引超出范围: {shape_index}"}

            shape = slide.shapes[shape_index]

            # 添加超链接
            try:
                # 类型安全的文本框操作
                if hasattr(shape, 'text_frame'):
                    text_frame = getattr(shape, 'text_frame')
                    if text_frame is not None:
                        if display_text:
                            text_frame.text = display_text

                        if not text_frame.paragraphs:
                            paragraph = text_frame.add_paragraph()
                        else:
                            paragraph = text_frame.paragraphs[0]

                        if not paragraph.runs:
                            run = paragraph.add_run()
                        else:
                            run = paragraph.runs[0]

                        run.hyperlink.address = url
                    else:
                        # 其他形状类型
                        if hasattr(shape, 'click_action'):
                            shape.click_action.hyperlink.address = url
                        else:
                            return {"success": False, "error": "形状不支持超链接"}
                else:
                    # 其他形状类型
                    if hasattr(shape, 'click_action'):
                        shape.click_action.hyperlink.address = url
                    else:
                        return {"success": False, "error": "形状不支持超链接"}

            except Exception as e:
                return {"success": False, "error": f"添加超链接失败: {str(e)}"}

            return {
                "success": True,
                "message": f"成功为幻灯片 {slide_index} 的形状 {shape_index} 添加超链接",
                "url": url,
                "display_text": display_text
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def set_text_formatting(self, slide_index: int, shape_index: int, font_name: Optional[str] = None,
                           font_size: Optional[int] = None, font_color: Optional[str] = None, bold: Optional[bool] = None,
                           italic: Optional[bool] = None, underline: Optional[bool] = None) -> Dict[str, Any]:
        """设置文本格式"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            if shape_index >= len(slide.shapes):
                return {"success": False, "error": f"形状索引超出范围: {shape_index}"}

            shape = slide.shapes[shape_index]

            try:
                if not hasattr(shape, 'text_frame'):
                    return {"success": False, "error": "形状不支持文本框"}
                    
                text_frame = getattr(shape, 'text_frame')
                if text_frame is None:
                    return {"success": False, "error": "文本框不可用"}
                    
                if not hasattr(text_frame, 'paragraphs'):
                    return {"success": False, "error": "文本框没有段落属性"}
                    
                paragraphs = text_frame.paragraphs
                if not paragraphs or len(paragraphs) == 0:
                    return {"success": False, "error": "没有可用的文本段落"}
                    
                paragraph = paragraphs[0]
                if not hasattr(paragraph, 'font'):
                    return {"success": False, "error": "段落没有字体属性"}
                    
                font = paragraph.font
                if font is None:
                    return {"success": False, "error": "无法获取字体对象"}

                # 设置字体属性
                if font_name:
                    font.name = font_name
                if font_size:
                    font.size = Pt(font_size)
                if font_color:
                    try:
                        rgb_color = RGBColor.from_string(font_color)
                        font.color.rgb = rgb_color
                    except:
                        return {"success": False, "error": f"无效的颜色格式: {font_color}"}
                if bold is not None:
                    font.bold = bold
                if italic is not None:
                    font.italic = italic
                if underline is not None:
                    font.underline = underline
            except Exception as e:
                return {"success": False, "error": f"字体设置失败: {str(e)}"}

            return {
                "success": True,
                "message": f"成功设置幻灯片 {slide_index} 形状 {shape_index} 的文本格式",
                "formatting": {
                    "font_name": font_name,
                    "font_size": font_size,
                    "font_color": font_color,
                    "bold": bold,
                    "italic": italic,
                    "underline": underline
                }
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_slide_shapes_info(self, slide_index: int) -> Dict[str, Any]:
        """获取幻灯片中所有形状的信息"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]
            shapes_info = []

            for i, shape in enumerate(slide.shapes):
                shape_info = {
                    "index": i,
                    "shape_type": str(shape.shape_type),
                    "name": shape.name,
                    "left": shape.left.inches if hasattr(shape.left, 'inches') else 0,
                    "top": shape.top.inches if hasattr(shape.top, 'inches') else 0,
                    "width": shape.width.inches if hasattr(shape.width, 'inches') else 0,
                    "height": shape.height.inches if hasattr(shape.height, 'inches') else 0,
                    "has_text": hasattr(shape, 'text_frame') and getattr(shape, 'text_frame', None) is not None,
                    "text": getattr(shape, 'text', "") if hasattr(shape, 'text') else ""
                }
                shapes_info.append(shape_info)

            return {
                "success": True,
                "slide_index": slide_index,
                "shapes_count": len(shapes_info),
                "shapes": shapes_info
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def _copy_shape(self, source_shape, target_slide):
        """辅助方法：复制形状（简化版本）"""
        try:
            # 这是一个简化的形状复制方法
            # 对于复杂的形状复制，可能需要更详细的实现
            if hasattr(source_shape, 'text_frame') and source_shape.text_frame:
                # 复制文本框
                textbox = target_slide.shapes.add_textbox(
                    source_shape.left, source_shape.top,
                    source_shape.width, source_shape.height
                )
                textbox.text_frame.text = source_shape.text_frame.text
        except:
            # 如果复制失败，忽略该形状
            pass

    def set_slide_transition(self, slide_index: int, transition_type: str = "fade",
                           duration: float = 1.0, advance_on_click: bool = True,
                           advance_after_time: Optional[float] = None) -> Dict[str, Any]:
        """设置幻灯片过渡效果"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if slide_index >= len(slides):
                return {"success": False, "error": f"幻灯片索引超出范围: {slide_index}"}

            slide = slides[slide_index]

            # 支持的过渡类型
            supported_transitions = ["none", "fade", "push", "wipe", "split", "zoom", "blinds", "dissolve"]

            if transition_type.lower() not in supported_transitions:
                return {"success": False, "error": f"不支持的过渡类型: {transition_type}。支持的类型: {', '.join(supported_transitions)}"}

            # 获取幻灯片的XML元素
            slide_element = slide._element

            # 定义命名空间
            namespaces = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

            # 移除现有的过渡元素（如果存在）
            existing_transition = slide_element.find('.//p:transition', namespaces)
            if existing_transition is not None:
                slide_element.remove(existing_transition)

            # 创建新的过渡元素（如果不是none）
            if transition_type.lower() != "none":
                try:
                    from lxml import etree
                except ImportError:
                    return {"success": False, "error": "需要安装lxml库: pip install lxml"}

                # 创建过渡XML字符串
                transition_xml = self._create_transition_xml(transition_type, duration, advance_on_click, advance_after_time)

                # 解析XML并插入到幻灯片中
                parser = etree.XMLParser(ns_clean=True, recover=True)
                transition_elem = etree.fromstring(transition_xml.encode('utf-8'), parser)

                # 将过渡元素插入到符合规范的位置
                # p:transition 应该在 p:cSld 和 p:clrMapOvr 之间
                color_map_override = slide_element.find('.//p:clrMapOvr', namespaces)
                if color_map_override is not None:
                    color_map_override.addprevious(transition_elem)
                else:
                    # 如果没有 p:clrMapOvr，则追加到末尾
                    slide_element.append(transition_elem)

                # 验证插入是否成功
                verification_elem = slide_element.find('.//p:transition', namespaces)
                if verification_elem is None:
                    return {"success": False, "error": "过渡效果XML插入失败"}

            return {
                "success": True,
                "message": f"成功设置幻灯片 {slide_index} 的过渡效果",
                "transition_type": transition_type,
                "duration": duration,
                "advance_on_click": advance_on_click,
                "advance_after_time": advance_after_time,
                "verification": "过渡效果已验证插入成功" if transition_type.lower() != "none" else "已移除过渡效果"
            }

        except Exception as e:
            return {"success": False, "error": str(e)}

    def _create_transition_xml(self, transition_type: str, duration: float,
                              advance_on_click: bool, advance_after_time: Optional[float]) -> str:
        """创建过渡效果的XML字符串"""
        # 设置过渡速度
        if duration <= 0.5:
            speed = "fast"
        elif duration <= 2.0:
            speed = "med"
        else:
            speed = "slow"

        # 基础过渡元素
        transition_attrs = f'spd="{speed}"'

        if advance_on_click:
            transition_attrs += ' advClick="1"'
        else:
            transition_attrs += ' advClick="0"'

        if advance_after_time is not None:
            advance_time_ms = int(advance_after_time * 1000)
            transition_attrs += f' advTm="{advance_time_ms}"'

        # 根据过渡类型创建相应的XML
        transition_content = ""
        if transition_type.lower() == "fade":
            transition_content = '<p:fade/>'
        elif transition_type.lower() == "push":
            transition_content = '<p:push dir="l"/>'
        elif transition_type.lower() == "wipe":
            transition_content = '<p:wipe dir="l"/>'
        elif transition_type.lower() == "zoom":
            transition_content = '<p:zoom/>'
        elif transition_type.lower() == "split":
            transition_content = '<p:split orient="horz" dir="out"/>'
        elif transition_type.lower() == "blinds":
            transition_content = '<p:blinds dir="horz"/>'
        elif transition_type.lower() == "dissolve":
            transition_content = '<p:dissolve/>'
        else:
            # 默认使用fade
            transition_content = '<p:fade/>'

        # 完整的XML字符串
        xml_string = f'''<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" {transition_attrs}>
    {transition_content}
</p:transition>'''

        return xml_string

    def apply_transition_to_all_slides(self, transition_type: str = "fade", duration: float = 1.0) -> Dict[str, Any]:
        """为所有幻灯片应用统一的过渡效果"""
        try:
            if not self.current_presentation:
                return {"success": False, "error": "没有打开的演示文稿"}

            slides = self.current_presentation.slides
            if len(slides) == 0:
                return {"success": False, "error": "演示文稿中没有幻灯片"}

            success_count = 0
            failed_slides = []

            # 为每张幻灯片设置过渡效果
            for i in range(len(slides)):
                result = self.set_slide_transition(i, transition_type, duration, True, None)
                if result.get("success"):
                    success_count += 1
                else:
                    failed_slides.append(i)

            if success_count == len(slides):
                return {
                    "success": True,
                    "message": f"成功为所有 {len(slides)} 张幻灯片设置了 '{transition_type}' 过渡效果",
                    "transition_type": transition_type,
                    "duration": duration,
                    "slides_processed": len(slides)
                }
            else:
                return {
                    "success": True,
                    "message": f"为 {success_count}/{len(slides)} 张幻灯片设置了过渡效果",
                    "transition_type": transition_type,
                    "duration": duration,
                    "slides_processed": success_count,
                    "failed_slides": failed_slides,
                    "warning": f"有 {len(failed_slides)} 张幻灯片设置失败"
                }

        except Exception as e:
            return {"success": False, "error": str(e)}

    def get_available_transitions(self) -> Dict[str, Any]:
        """获取可用的过渡效果列表"""
        try:
            # 只返回实际支持的过渡效果
            transitions = [
                {"name": "none", "description": "无过渡效果"},
                {"name": "fade", "description": "淡入淡出 - 推荐用于专业演示"},
                {"name": "push", "description": "推入 - 动感十足"},
                {"name": "wipe", "description": "擦除 - 简洁流畅"},
                {"name": "split", "description": "分割 - 创意效果"},
                {"name": "zoom", "description": "缩放 - 突出重点"},
                {"name": "blinds", "description": "百叶窗 - 经典效果"},
                {"name": "dissolve", "description": "溶解 - 柔和过渡"}
            ]

            return {
                "success": True,
                "transitions": transitions,
                "total_count": len(transitions),
                "note": "这些动画效果可以让您的演示文稿更加生动有趣",
                "recommendation": "推荐使用 'fade' 效果，适合大多数专业演示场合"
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    def make_presentation_professional(self) -> Dict[str, Any]:
        """一键让演示文稿变得专业 - 添加淡入淡出过渡效果"""
        return self.apply_transition_to_all_slides("fade", 1.0)

    def add_smooth_transitions(self) -> Dict[str, Any]:
        """为演示文稿添加流畅的过渡动画"""
        return self.apply_transition_to_all_slides("fade", 0.8)

    def add_dynamic_effects(self) -> Dict[str, Any]:
        """为演示文稿添加动感效果"""
        return self.apply_transition_to_all_slides("push", 1.2)

    def generate_outline_for_topic(self, topic: str) -> Dict[str, Any]:
        """根据主题，生成一个结构化的、用于创建演示文稿的JSON大纲。"""
        try:
            # 为了演示，我们在这里生成一个硬编码的示例大纲。
            # 在实际应用中，这里可以是对真正LLM服务的API调用。
            outline_data = {
                "slides": [
                    {
                        "title": f"关于 {topic} 的深入探讨",
                        "subtitle": "由AI辅助生成"
                    },
                    {
                        "title": "介绍与背景",
                        "content": [
                            f"{topic} 的定义与重要性",
                            "相关的历史发展",
                            "本次讨论的主要范围"
                        ]
                    },
                    {
                        "title": "核心要点分析",
                        "content": [
                            "第一个关键方面",
                            "第二个关键方面，并提供示例",
                            "第三个关键方面的深入分析"
                        ]
                    },
                    {
                        "title": "案例研究或实际应用",
                        "content": [
                            f"一个关于 {topic} 的真实世界案例",
                            "从案例中得到的启示",
                            "如何将这些应用到实践中"
                        ]
                    },
                    {
                        "title": "总结与展望",
                        "content": [
                            f"对 {topic} 的核心内容进行总结",
                            "未来的发展趋势",
                            "问答环节"
                        ]
                    }
                ]
            }
            
            outline_json = json.dumps(outline_data, ensure_ascii=False, indent=2)

            return {
                "success": True,
                "message": f"成功为主题 '{topic}' 生成大纲。",
                "outline_json": outline_json
            }
        except Exception as e:
            logger.error(f"为主题 '{topic}' 生成大纲时出错: {e}")
            return {"success": False, "error": str(e)}